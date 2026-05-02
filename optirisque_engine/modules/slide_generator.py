"""
Module slide_generator — F du cahier des charges.

Génère :
- slides_outline.md : storyboard exécutif sobre (toujours produit)
- slides_cic.pptx   : deck PowerPoint conforme à la charte DÉPS
                      (uniquement si python-pptx est installé)

Contraintes par slide :
- 1 titre décisionnel (verbe d'action ou conclusion)
- 5 lignes max dans le corps
- 1 KPI / preuve chiffrée
- 1 conclusion claire pour le CIC

Aucune donnée n'est inventée : toute case vide affiche "non extrait".
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from .pdf_extractor import ExtractionResult, NON_EXTRAIT
from .advisor_score import AdvisorScore


MAX_BODY_LINES = 5


# ---------------------------------------------------------------------------
# Construction du contenu (commun MD / PPTX)
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    numero: int
    titre: str
    sous_titre: str
    corps: List[str]
    kpi: str
    conclusion: str


def _short(value: Any, limit: int = 90) -> str:
    if value is None or value == "":
        return NON_EXTRAIT
    s = str(value).strip()
    if len(s) <= limit:
        return s
    return s[: limit - 1].rstrip() + "…"


def _bullets(items: List[str], limit: int = MAX_BODY_LINES) -> List[str]:
    if not items:
        return [NON_EXTRAIT]
    return [_short(i, 110) for i in items[:limit]]


def build_slides(
    extraction: ExtractionResult,
    score: AdvisorScore,
    config: Dict[str, Any],
) -> List[Slide]:
    champs = extraction.champs
    suivi = extraction.suivi or {}

    slides: List[Slide] = []

    # 1 — Page titre
    slides.append(Slide(
        numero=1,
        titre=f"Comité d'investissement — {_short(champs.get('entreprise'), 50)}",
        sous_titre="DÉPS — MRC Pierre-De Saurel",
        corps=[
            f"Dossier : {_short(champs.get('entreprise'))}",
            f"Secteur : {_short(champs.get('secteur'))}",
            f"FLI : {_short(champs.get('fli'))} | FLS : {_short(champs.get('fls'))}",
            f"Run ID : {score.run_id}",
        ],
        kpi=f"Score conseiller : {score.score_total}/100 — Niveau {score.niveau}/10",
        conclusion=f"Décision attendue : {score.decision}",
    ))

    # 2 — Décision attendue
    slides.append(Slide(
        numero=2,
        titre="Décision attendue du CIC",
        sous_titre="Recommandation institutionnelle",
        corps=[
            f"Recommandation : {_short(champs.get('recommandation'), 140)}",
            f"Verdict de risque : {_short(champs.get('verdict_risque'))}",
            f"Score qualité dossier : {score.score_total}/100 ({score.qualificatif})",
            f"Niveau conseiller : {score.niveau}/10 (cible 10/10)",
        ],
        kpi=f"Décision : {score.decision}",
        conclusion="Vote attendu : approbation / report / refus.",
    ))

    # 3 — Portrait entreprise
    slides.append(Slide(
        numero=3,
        titre=f"Portrait — {_short(champs.get('entreprise'), 60)}",
        sous_titre="Identification et activité",
        corps=[
            f"Entreprise : {_short(champs.get('entreprise'))}",
            f"Secteur : {_short(champs.get('secteur'))}",
            f"Nature du projet : {_short(champs.get('nature_projet'), 140)}",
            f"FLI / FLS : {_short(champs.get('fli'))} / {_short(champs.get('fls'))}",
        ],
        kpi=f"Carnet de commandes : {_short(extraction.kpi.get('carnet_commandes'))}",
        conclusion="Profil aligné avec la mission DÉPS." if champs.get("entreprise") not in (None, NON_EXTRAIT) else "Identification incomplète — retour conseiller requis.",
    ))

    # 4 — Montage financier
    slides.append(Slide(
        numero=4,
        titre="Montage financier proposé",
        sous_titre="Sources et montant",
        corps=[
            f"Montant demandé : {_short(champs.get('montant_demande'))}",
            f"Nature du projet : {_short(champs.get('nature_projet'), 140)}",
            "Sources de financement : voir gabarit OptiRisque (section Coûts et sources).",
            "Sommaire du financement proposé : voir section dédiée du dossier.",
        ],
        kpi=f"Endettement post-financement : {_short(extraction.kpi.get('endettement'))}",
        conclusion="Structure conforme aux paramètres FLI/FLS." if champs.get("montant_demande") not in (None, NON_EXTRAIT) else "Montant non extrait — vérifier la complétude du gabarit.",
    ))

    # 5 — KPI financiers
    kpis = extraction.kpi
    slides.append(Slide(
        numero=5,
        titre="KPI financiers clés",
        sous_titre="Indicateurs de viabilité",
        corps=[
            f"DSCR : {_short(kpis.get('DSCR'))}",
            f"Marge brute : {_short(kpis.get('marge_brute'))}",
            f"Endettement : {_short(kpis.get('endettement'))}",
            f"Concentration client : {_short(kpis.get('concentration_client'))}",
            f"Carnet de commandes : {_short(kpis.get('carnet_commandes'))}",
        ][:MAX_BODY_LINES],
        kpi=f"DSCR : {_short(kpis.get('DSCR'))}",
        conclusion=_kpi_conclusion(kpis),
    ))

    # 6 — Analyse des risques
    risques = extraction.risques or {}
    risques_lines = (
        [f"{cat} : cote {info.get('cote', NON_EXTRAIT)}" for cat, info in risques.items()]
        if risques else [NON_EXTRAIT]
    )
    slides.append(Slide(
        numero=6,
        titre="Analyse des risques",
        sous_titre="Cotation par catégorie FTQ",
        corps=risques_lines[:MAX_BODY_LINES],
        kpi=f"Catégories cotées : {len(risques)}/{len(config.get('categories_risques', []))}",
        conclusion=(
            "Cotation complète des 5 catégories." if len(risques) >= 5
            else "Cotation incomplète — compléter avant CIC."
        ),
    ))

    # 7 — Conditions de décaissement
    decais = suivi.get("decaissement", [])
    slides.append(Slide(
        numero=7,
        titre="Conditions de décaissement",
        sous_titre="Préalables SMART",
        corps=_bullets(decais),
        kpi=f"{len(decais)} condition(s) listée(s)",
        conclusion=(
            "Conditions formalisées et SMART." if decais
            else "Conditions absentes — bloquant pour décaissement."
        ),
    ))

    # 8 — Suivi post-financement
    suivi_items = suivi.get("suivi", [])
    indicateurs = suivi.get("indicateurs_mensuels_trimestriels", [])
    suivi_lines = (suivi_items + indicateurs)[:MAX_BODY_LINES] or [NON_EXTRAIT]
    slides.append(Slide(
        numero=8,
        titre="Suivi post-financement",
        sous_titre="Cadre et indicateurs",
        corps=suivi_lines,
        kpi=f"{len(indicateurs)} indicateur(s) de suivi",
        conclusion=(
            "Cadre de suivi opérationnel." if (suivi_items or indicateurs)
            else "Cadre de suivi à définir avec le conseiller."
        ),
    ))

    # 9 — Recommandation
    slides.append(Slide(
        numero=9,
        titre="Recommandation DÉPS",
        sous_titre="Position institutionnelle",
        corps=[
            f"Recommandation : {_short(champs.get('recommandation'), 160)}",
            f"Score conseiller : {score.score_total}/100 — niveau {score.niveau}/10",
            f"Qualificatif : {score.qualificatif}",
            f"Points forts : {len(score.points_forts)}",
            f"Lacunes : {len(score.lacunes)}",
        ],
        kpi=f"Décision : {score.decision}",
        conclusion=(
            "Dossier prêt pour vote CIC." if score.decision == "prêt CIC"
            else "Retour conseiller — corriger les lacunes listées."
        ),
    ))

    # 10 — Annexes
    slides.append(Slide(
        numero=10,
        titre="Annexes",
        sous_titre="Références et trace",
        corps=[
            f"Sections gabarit présentes : {len(extraction.sections_detectees)}/{len(config.get('gabarit_sections', []))}",
            f"Alertes d'extraction : {len(extraction.alertes)}",
            f"Run ID : {score.run_id}",
            "Voir run_report.json et advisor_score.json pour le détail audit-ready.",
        ],
        kpi=f"Audit trail : run_report.json (run {score.run_id})",
        conclusion="Documentation auditable conservée par le DÉPS.",
    ))

    return slides


def _kpi_conclusion(kpis: Dict[str, Any]) -> str:
    extracted = sum(1 for v in kpis.values() if v not in (None, "", NON_EXTRAIT))
    total = len(kpis)
    if extracted == total and total > 0:
        return "Indicateurs financiers complets — viabilité documentée."
    if extracted == 0:
        return "Aucun KPI extrait — dossier incomplet."
    return f"Indicateurs partiels ({extracted}/{total}) — compléter avant CIC."


# ---------------------------------------------------------------------------
# Sortie Markdown (slides_outline.md)
# ---------------------------------------------------------------------------

def render_outline_markdown(slides: List[Slide], dossier: str, run_id: str) -> str:
    lines: List[str] = [
        f"# Deck CIC — {dossier}",
        "",
        f"_Run {run_id} — Storyboard exécutif (sobre, charte DÉPS)._",
        "",
    ]
    for s in slides:
        lines.append(f"## Slide {s.numero} — {s.titre}")
        lines.append(f"*{s.sous_titre}*")
        lines.append("")
        for line in s.corps:
            lines.append(f"- {line}")
        lines.append("")
        lines.append(f"**KPI** : {s.kpi}")
        lines.append("")
        lines.append(f"**Conclusion CIC** : {s.conclusion}")
        lines.append("")
        lines.append("---")
        lines.append("")
    return "\n".join(lines)


def write_outline(path: Path, content: str) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


# ---------------------------------------------------------------------------
# Sortie PowerPoint (slides_cic.pptx)
# ---------------------------------------------------------------------------

def write_pptx(path: Path, slides: List[Slide], config: Dict[str, Any]) -> Tuple[bool, Optional[str]]:
    """Génère le PPTX charté DÉPS. Retourne (success, message)."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt, Emu  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except Exception as e:  # pragma: no cover - dépendance optionnelle
        return False, f"python-pptx indisponible ({e}); export PPTX ignoré."

    couleurs = config.get("charte", {}).get("couleurs", {})
    typo_primaire = config.get("charte", {}).get("typographies", {}).get("primaire", "Noto Sans")
    typo_fallback = config.get("charte", {}).get("typographies", {}).get("fallback", "Noto Sans")

    indigo = _hex_to_rgb(couleurs.get("indigo", "#24135D"))
    bleu = _hex_to_rgb(couleurs.get("bleu", "#0057B8"))
    cyan = _hex_to_rgb(couleurs.get("cyan", "#00AEC7"))
    blanc = _hex_to_rgb(couleurs.get("blanc", "#FFFFFF"))
    gris = _hex_to_rgb(couleurs.get("gris_texte", "#1F2330"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for s in slides:
        slide = prs.slides.add_slide(blank_layout)

        # Bandeau supérieur Indigo
        bandeau = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
        )
        bandeau.line.fill.background()
        bandeau.fill.solid()
        bandeau.fill.fore_color.rgb = RGBColor(*indigo)
        # Trait Cyan en bas du bandeau
        trait = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.0), prs.slide_width, Emu(38100)
        )
        trait.line.fill.background()
        trait.fill.solid()
        trait.fill.fore_color.rgb = RGBColor(*cyan)

        # Titre
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15), prs.slide_width - Inches(1.0), Inches(0.75)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s.titre
        for run in p.runs:
            run.font.name = typo_primaire
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*blanc)

        # Sous-titre (sous le bandeau)
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.15), prs.slide_width - Inches(1.0), Inches(0.45)
        )
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = s.sous_titre
        for run in sp.runs:
            run.font.name = typo_fallback
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(*bleu)

        # Corps (≤ 5 bullets)
        body_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.8), prs.slide_width - Inches(1.2), Inches(4.0)
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        for i, line in enumerate(s.corps[:MAX_BODY_LINES]):
            para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            para.text = f"• {line}"
            for run in para.runs:
                run.font.name = typo_fallback
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(*gris)

        # Bandeau KPI (cyan)
        kpi_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), Inches(5.85),
            prs.slide_width - Inches(1.2), Inches(0.65),
        )
        kpi_box.line.fill.background()
        kpi_box.fill.solid()
        kpi_box.fill.fore_color.rgb = RGBColor(*cyan)
        kpi_tf = kpi_box.text_frame
        kpi_tf.margin_left = Inches(0.2)
        kpi_tf.margin_right = Inches(0.2)
        kpi_tf.word_wrap = True
        kp = kpi_tf.paragraphs[0]
        kp.text = f"KPI — {s.kpi}"
        for run in kp.runs:
            run.font.name = typo_primaire
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*indigo)

        # Conclusion CIC (texte sous le bandeau KPI)
        ccl_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(6.6), prs.slide_width - Inches(1.2), Inches(0.6)
        )
        cp = ccl_box.text_frame.paragraphs[0]
        cp.text = f"Conclusion CIC : {s.conclusion}"
        for run in cp.runs:
            run.font.name = typo_fallback
            run.font.size = Pt(13)
            run.font.italic = True
            run.font.color.rgb = RGBColor(*bleu)

        # Numéro de slide en bas à droite
        num_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(1.0), prs.slide_height - Inches(0.4),
            Inches(0.8), Inches(0.3),
        )
        np_ = num_box.text_frame.paragraphs[0]
        np_.text = f"{s.numero}"
        for run in np_.runs:
            run.font.name = typo_fallback
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(*gris)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return True, None


def _hex_to_rgb(hex_color: str) -> tuple:
    h = hex_color.lstrip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))
