from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# ---------- STYLE GLOBAL ----------
def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 25, 47)  # Bleu nuit

def add_title(slide, text):
    title = slide.shapes.title
    title.text = text
    for p in title.text_frame.paragraphs:
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

def add_kpi_box(slide, left, top, width, height, title, value):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 174, 199)  # turquoise DÉPS

    text_frame = shape.text_frame
    text_frame.clear()

    p1 = text_frame.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = title
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER


# ---------- SLIDE 1 ----------
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)

add_title(slide, "Pilotage stratégique du portefeuille FLI / FLS")

subtitle = slide.placeholders[1]
subtitle.text = "Performance, gestion des risques et impact économique\nDÉPS"


# ---------- SLIDE 2 (KPI VISUEL) ----------
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)

add_title(slide, "Indicateurs clés du portefeuille")

add_kpi_box(slide, Inches(1), Inches(2), Inches(3), Inches(2),
            "IRR pondéré", "8,4 %")

add_kpi_box(slide, Inches(4.5), Inches(2), Inches(3), Inches(2),
            "Taux de défaut", "3,2 %")

add_kpi_box(slide, Inches(8), Inches(2), Inches(3), Inches(2),
            "Effet levier", "2,7x")


# ---------- SLIDE 3 ----------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_dark_background(slide)

add_title(slide, "Structure du portefeuille")

tf = slide.placeholders[1].text_frame
tf.text = "Encours : 12,4 M$"
p = tf.add_paragraph()
p.text = "86 dossiers actifs"
p = tf.add_paragraph()
p.text = "Ticket moyen : 144 000 $"
p = tf.add_paragraph()
p.text = "Croissance : +9 %"


# ---------- SLIDE 4 ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_dark_background(slide)

add_title(slide, "Répartition sectorielle")

tf = slide.placeholders[1].text_frame
tf.text = "Manufacturier : 28 %"
for line in ["Services : 22 %", "Commerce : 18 %",
             "Agroalimentaire : 16 %", "Construction : 10 %", "Autres : 6 %"]:
    p = tf.add_paragraph()
    p.text = line


# ---------- SLIDE 5 ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_dark_background(slide)

add_title(slide, "Profil de risque")

tf = slide.placeholders[1].text_frame
tf.text = "Faible : 42 %"
for line in ["Modéré : 38 %", "Élevé : 15 %", "Critique : 5 %"]:
    p = tf.add_paragraph()
    p.text = line


# ---------- SLIDE 6 ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_dark_background(slide)

add_title(slide, "Performance financière")

tf = slide.placeholders[1].text_frame
tf.text = "2023 : 5,2 % / 4,1 % défaut"
for line in ["2024 : 6,0 % / 3,8 %",
             "2025 : 6,1 % / 3,2 %"]:
    p = tf.add_paragraph()
    p.text = line


# ---------- SLIDE 7 ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_dark_background(slide)

add_title(slide, "Impact économique")

tf = slide.placeholders[1].text_frame
tf.text = "33,5 M$ générés"
for line in ["Effet levier 2,7x",
             "412 emplois",
             "18 projets relève"]:
    p = tf.add_paragraph()
    p.text = line


# ---------- SLIDE 8 ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_dark_background(slide)

add_title(slide, "Enjeux stratégiques")

tf = slide.placeholders[1].text_frame
tf.text = "Pression sectorielle"
for line in ["Hausse des taux",
             "Concentration dossiers",
             "Risque à surveiller"]:
    p = tf.add_paragraph()
    p.text = line


# ---------- SLIDE 9 ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_dark_background(slide)

add_title(slide, "Recommandations")

tf = slide.placeholders[1].text_frame
tf.text = "Prioriser projets à forte valeur"
for line in ["Renforcer analyse des risques",
             "Accroître accompagnement",
             "Maintenir diversification"]:
    p = tf.add_paragraph()
    p.text = line


# ---------- SLIDE 10 ----------
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_dark_background(slide)

add_title(slide, "Conclusion")

tf = slide.placeholders[1].text_frame
tf.text = "Portefeuille solide"
for line in ["Croissance contrôlée",
             "Impact durable"]:
    p = tf.add_paragraph()
    p.text = line


# ---------- EXPORT ----------
import os
output_dir = os.path.dirname(os.path.abspath(__file__))
file_path = os.path.join(output_dir, "DEPS_Executive_Premium.pptx")
prs.save(file_path)

print(file_path)
